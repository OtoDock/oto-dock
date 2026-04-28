"""write_pptx add_chart color styling — series colors + pie slice colors.

Needs python-pptx; skips where it's absent (the MCP's venv / image CI).
"""

import asyncio
import sys
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")

sys.path.insert(0, str(Path(__file__).parent.parent))

import powerpoint as ppt_mod  # noqa: E402


@pytest.fixture(autouse=True)
def _bypass_platform(monkeypatch):
    monkeypatch.setattr(ppt_mod, "_resolve_path", lambda p, writing=False: p)
    monkeypatch.setattr(ppt_mod, "_to_agents_relative", lambda p: p)

    async def _noop_preview(path):
        pass

    monkeypatch.setattr(ppt_mod, "_push_preview", _noop_preview)


def _write(path: Path, ops: list[dict]) -> str:
    return asyncio.run(ppt_mod.handle_write_pptx(
        {"path": str(path), "create_new": True, "operations": ops}))


def test_add_chart_series_color_applied(tmp_path):
    f = tmp_path / "deck.pptx"
    out = _write(f, [
        {"type": "add_slide", "layout": "blank"},
        {"type": "add_chart", "slide": 0, "chart_type": "column_clustered",
         "categories": ["Q1", "Q2"],
         "series": [{"name": "Rev", "values": [1, 2], "color": "#16BAC5"}]},
    ])
    assert "Warnings/Errors" not in out, out

    from pptx import Presentation
    prs = Presentation(str(f))
    chart = next(s for s in prs.slides[0].shapes if s.has_chart).chart
    ser = list(chart.series)[0]
    assert str(ser.format.fill.fore_color.rgb) == "16BAC5"


def test_add_chart_pie_slice_colors_applied(tmp_path):
    f = tmp_path / "pie.pptx"
    out = _write(f, [
        {"type": "add_slide", "layout": "blank"},
        {"type": "add_chart", "slide": 0, "chart_type": "doughnut",
         "categories": ["a", "b", "c"],
         "series": [{"name": "S", "values": [5, 3, 2]}],
         "colors": ["#111111", "#222222", "#333333"]},
    ])
    assert "Warnings/Errors" not in out, out

    from pptx import Presentation
    prs = Presentation(str(f))
    chart = next(s for s in prs.slides[0].shapes if s.has_chart).chart
    points = list(list(chart.series)[0].points)
    got = [str(p.format.fill.fore_color.rgb) for p in points]
    assert got == ["111111", "222222", "333333"]