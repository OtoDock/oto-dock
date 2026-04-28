"""PowerPoint (PPTX) read and write handlers for the file-tools MCP.

Enhanced reader with shape inventory, chart data, speaker notes.
26 write operations including shapes, charts, speaker notes, table cell
formatting, transitions, connectors, hyperlinks, and slide management.
"""

import copy
from pathlib import Path

from shared import _normalize_operations, _op_type, _push_preview, _resolve_path, _to_agents_relative, logger

# ---------------------------------------------------------------------------
# Shape type mapping
# ---------------------------------------------------------------------------

_SHAPE_TYPE_MAP = None


def _get_shape_map():
    global _SHAPE_TYPE_MAP
    if _SHAPE_TYPE_MAP is None:
        from pptx.enum.shapes import MSO_SHAPE
        _SHAPE_TYPE_MAP = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
            "arrow_up": MSO_SHAPE.UP_ARROW,
            "arrow_down": MSO_SHAPE.DOWN_ARROW,
            "star": MSO_SHAPE.STAR_5_POINT,
            "star_5": MSO_SHAPE.STAR_5_POINT,
            "star_6": MSO_SHAPE.STAR_6_POINT,
            "star_4": MSO_SHAPE.STAR_4_POINT,
            "chevron": MSO_SHAPE.CHEVRON,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "cloud": MSO_SHAPE.CLOUD,
            "heart": MSO_SHAPE.HEART,
            "lightning": MSO_SHAPE.LIGHTNING_BOLT,
            "plus": MSO_SHAPE.CROSS,
            "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
            "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
            "flowchart_data": MSO_SHAPE.FLOWCHART_DATA,
            "flowchart_terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
            "callout_1": MSO_SHAPE.RECTANGULAR_CALLOUT,
            "callout_2": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        }
    return _SHAPE_TYPE_MAP


# ---------------------------------------------------------------------------
# Chart type mapping
# ---------------------------------------------------------------------------

_CHART_TYPE_MAP = None


def _get_chart_type_map():
    global _CHART_TYPE_MAP
    if _CHART_TYPE_MAP is None:
        from pptx.enum.chart import XL_CHART_TYPE
        _CHART_TYPE_MAP = {
            "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "scatter_lines": XL_CHART_TYPE.XY_SCATTER_LINES,
            "area": XL_CHART_TYPE.AREA,
            "area_stacked": XL_CHART_TYPE.AREA_STACKED,
            "radar": XL_CHART_TYPE.RADAR,
        }
    return _CHART_TYPE_MAP


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _parse_rgb(color_str: str):
    from pptx.dml.color import RGBColor
    c = str(color_str).lstrip("#")
    return RGBColor(int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _apply_run_format(run, fmt: dict):
    """Apply formatting to a pptx run."""
    from pptx.util import Pt
    if fmt.get("font_size"):
        run.font.size = Pt(int(fmt["font_size"]))
    if fmt.get("bold") is not None:
        run.font.bold = fmt["bold"]
    if fmt.get("italic") is not None:
        run.font.italic = fmt["italic"]
    if fmt.get("underline"):
        run.font.underline = True
    color = fmt.get("color") or fmt.get("font_color") or fmt.get("text_color")
    if color:
        run.font.color.rgb = _parse_rgb(color)
    if fmt.get("font_name"):
        run.font.name = fmt["font_name"]


def _get_slide(prs, op: dict):
    """Get slide by index from operation. Supports negative indexing (-1 = last slide)."""
    raw = op.get("slide_index")
    if raw is None:
        raw = op.get("slide")
    if raw is None:
        raise IndexError("slide_index is required")
    si = int(raw)
    n = len(prs.slides)
    # Support negative indexing: -1 = last slide, -2 = second to last, etc.
    if si < 0:
        si = n + si
    if si < 0 or si >= n:
        raise IndexError(f"slide_index {raw} out of range (have {n} slides, use 0-{n-1} or -1 for last)")
    return prs.slides[si]


def _find_shape(slide, op: dict):
    """Find a shape by name, text content, or index on a slide."""
    name = op.get("shape_name")
    if name:
        # Try exact shape.name match first
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        # Fall back to matching text content
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text == name:
                return shape
        raise ValueError(f"Shape '{name}' not found on slide")
    idx = op.get("shape_index")
    if idx is not None:
        idx = int(idx)
        shapes = list(slide.shapes)
        if 0 <= idx < len(shapes):
            return shapes[idx]
        raise IndexError(f"shape_index {idx} out of range")
    raise ValueError("Must provide shape_name or shape_index")


# ---------------------------------------------------------------------------
# Read
# ---------------------------------------------------------------------------


def read_pptx(path: str) -> str:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(path)
    total = len(prs.slides)
    w_in = round(prs.slide_width / Emu(914400), 2) if prs.slide_width else "?"
    h_in = round(prs.slide_height / Emu(914400), 2) if prs.slide_height else "?"
    result = [f"**PPTX**: {Path(path).name} — {total} slide(s), {w_in}x{h_in}in"]

    # Properties
    try:
        props = prs.core_properties
        meta = []
        if props.title:
            meta.append(f"Title: {props.title}")
        if props.author:
            meta.append(f"Author: {props.author}")
        if meta:
            result.append("**Properties**: " + " | ".join(meta))
    except Exception:
        pass

    result.append("")

    for i, slide in enumerate(prs.slides):
        layout = slide.slide_layout.name if slide.slide_layout else "unknown"
        result.append(f"--- Slide {i+1} (layout: {layout}) ---")

        # Speaker notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    result.append(f"  Notes: {notes}")
        except Exception:
            pass

        # Shape inventory
        for shape in slide.shapes:
            prefix = f"  [{shape.shape_type}] {shape.name}"
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    result.append(f"{prefix}: {text[:100]}{'...' if len(text) > 100 else ''}")
                else:
                    result.append(prefix)
            elif shape.shape_type == 13:  # Picture
                result.append(f"  [Picture] {shape.name}")
            elif hasattr(shape, "has_chart") and shape.has_chart:
                chart = shape.chart
                ct = str(chart.chart_type) if chart.chart_type else "unknown"
                title = chart.chart_title.text_frame.text if chart.has_title else ""
                result.append(f"  [Chart] {shape.name}: {ct}" + (f' "{title}"' if title else ""))
            elif hasattr(shape, "has_table") and shape.has_table:
                t = shape.table
                result.append(f"  [Table] {shape.name}: {len(t.rows)}x{len(t.columns)}")
            else:
                result.append(prefix)

        result.append("")
    return "\n".join(result)


# ---------------------------------------------------------------------------
# Write — main handler
# ---------------------------------------------------------------------------


async def handle_write_pptx(args: dict) -> str:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    path = _resolve_path(args["path"], writing=True)
    ops = _normalize_operations(args.get("operations"))
    create_new = args.get("create_new", False)

    if Path(path).exists() and not create_new:
        prs = Presentation(path)
    else:
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()

    layout_map = {
        "title": 0, "title_and_content": 1, "section_header": 2,
        "two_content": 3, "comparison": 4, "title_only": 5,
        "blank": 6, "content": 1, "image": 6,
    }

    align_map = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY,
    }

    errors = []

    for idx, op in enumerate(ops):
        ot = _op_type(op)
        try:
            # =============================================================
            # SLIDES
            # =============================================================

            if ot == "add_slide":
                layout_name = op.get("layout", "blank")
                li = layout_map.get(layout_name, 6)
                li = min(li, len(prs.slide_layouts) - 1)
                prs.slides.add_slide(prs.slide_layouts[li])

            elif ot == "delete_slide":
                _si = op.get("slide_index")
                if _si is None:
                    _si = op.get("slide")
                si = int(_si if _si is not None else -1)
                n = len(prs.slides)
                if si < 0:
                    si = n + si
                if 0 <= si < n:
                    rId = prs.slides._sldIdLst[si].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[si]

            elif ot == "duplicate_slide":
                _si = op.get("slide_index")
                if _si is None:
                    _si = op.get("slide")
                si = int(_si if _si is not None else -1)
                if si < 0:
                    si = len(prs.slides) + si
                if si >= len(prs.slides):
                    errors.append(f"Op #{idx} duplicate_slide: index {si} out of range")
                    continue
                src = prs.slides[si]
                layout = src.slide_layout
                new_slide = prs.slides.add_slide(layout)
                # Copy all shapes from source to new slide
                for shape in src.shapes:
                    el = copy.deepcopy(shape._element)
                    new_slide.shapes._spTree.append(el)
                # Copy notes
                try:
                    if src.has_notes_slide:
                        notes_text = src.notes_slide.notes_text_frame.text
                        if notes_text:
                            new_slide.notes_slide.notes_text_frame.text = notes_text
                except Exception:
                    pass

            elif ot == "move_slide":
                from_idx = int(op.get("from_index", 0))
                to_idx = int(op.get("to_index", 0))
                sld_lst = prs.slides._sldIdLst
                if 0 <= from_idx < len(sld_lst) and 0 <= to_idx < len(sld_lst):
                    el = sld_lst[from_idx]
                    sld_lst.remove(el)
                    sld_lst.insert(to_idx, el)

            elif ot == "set_slide_dimensions":
                prs.slide_width = Inches(float(op.get("width", 13.333)))
                prs.slide_height = Inches(float(op.get("height", 7.5)))

            # =============================================================
            # TEXT & PLACEHOLDERS
            # =============================================================

            elif ot == "set_text":
                slide = _get_slide(prs, op)
                placeholder = op.get("placeholder", "title")
                ph_idx = int(placeholder) if placeholder.isdigit() else (0 if placeholder == "title" else 1)

                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == ph_idx:
                        # Support runs array for mixed formatting
                        runs_data = op.get("runs")
                        if runs_data and isinstance(runs_data, list):
                            shape.text = ""
                            tf = shape.text_frame
                            p = tf.paragraphs[0]
                            for rd in runs_data:
                                run = p.add_run()
                                run.text = rd.get("text", "")
                                _apply_run_format(run, rd)
                        else:
                            shape.text = op.get("text", "")
                            if shape.text_frame.paragraphs:
                                para = shape.text_frame.paragraphs[0]
                                run = para.runs[0] if para.runs else para.add_run()
                                if not para.runs:
                                    run.text = op.get("text", "")
                                _apply_run_format(run, op)
                        break

            elif ot == "add_textbox":
                slide = _get_slide(prs, op)
                left = Inches(float(op.get("left", 1)))
                top = Inches(float(op.get("top", 1)))
                width = Inches(float(op.get("width", 6)))
                height = Inches(float(op.get("height", 1)))
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = op.get("word_wrap", True)
                if op.get("auto_size"):
                    from pptx.enum.text import MSO_AUTO_SIZE
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                if op.get("margin") is not None:
                    m = Inches(float(op["margin"]))
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = m

                runs_data = op.get("runs")
                p = tf.paragraphs[0]
                if runs_data and isinstance(runs_data, list):
                    for rd in runs_data:
                        run = p.add_run()
                        run.text = rd.get("text", "")
                        _apply_run_format(run, rd)
                else:
                    run = p.add_run()
                    run.text = op.get("text", "")
                    _apply_run_format(run, op)

                alignment = op.get("alignment")
                if alignment and alignment in align_map:
                    p.alignment = align_map[alignment]

            elif ot == "add_bullet_points":
                slide = _get_slide(prs, op)
                # Find target text frame
                shape_name = op.get("shape_name")
                shape_idx = op.get("shape_index")
                ph_raw = op.get("placeholder")
                tf = None
                if shape_name:
                    # Match by shape.name then by text content
                    for s in slide.shapes:
                        if s.has_text_frame and (s.name == shape_name or s.text_frame.text == shape_name):
                            tf = s.text_frame
                            break
                elif shape_idx is not None:
                    shapes = list(slide.shapes)
                    si = int(shape_idx)
                    if 0 <= si < len(shapes) and shapes[si].has_text_frame:
                        tf = shapes[si].text_frame
                elif ph_raw is not None:
                    # Accept named placeholders or numeric index
                    ph_names = {"title": 0, "body": 1, "subtitle": 1, "content": 1}
                    ph_idx = ph_names.get(str(ph_raw).lower(), None)
                    if ph_idx is None:
                        try:
                            ph_idx = int(ph_raw)
                        except (ValueError, TypeError):
                            errors.append(f"Op #{idx} add_bullet_points: unknown placeholder '{ph_raw}'")
                            continue
                    for s in slide.placeholders:
                        if s.placeholder_format.idx == ph_idx:
                            tf = s.text_frame
                            break
                else:
                    # Default: body placeholder (idx 1)
                    for s in slide.placeholders:
                        if s.placeholder_format.idx == 1:
                            tf = s.text_frame
                            break

                if not tf:
                    errors.append(f"Op #{idx} add_bullet_points: no text frame found")
                    continue

                items = op.get("items", [])
                for item in items:
                    if isinstance(item, dict):
                        text = item.get("text", "")
                        level = int(item.get("level", 0))
                    else:
                        text = str(item)
                        level = 0
                    p = tf.add_paragraph()
                    p.text = text
                    p.level = min(level, 8)
                    if op.get("font_size"):
                        for run in p.runs:
                            run.font.size = Pt(int(op["font_size"]))

            elif ot == "format_text":
                slide = _get_slide(prs, op)
                find_text = op.get("find", "")
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if find_text in run.text:
                                _apply_run_format(run, op)

            # =============================================================
            # SHAPES
            # =============================================================

            elif ot == "add_shape":
                slide = _get_slide(prs, op)
                shape_type_name = op.get("shape_type", "rectangle")
                shape_map = _get_shape_map()
                mso_shape = shape_map.get(shape_type_name)
                if not mso_shape:
                    errors.append(f"Op #{idx} add_shape: unknown shape_type '{shape_type_name}'. Available: {', '.join(shape_map.keys())}")
                    continue

                left = Inches(float(op.get("left", 1)))
                top = Inches(float(op.get("top", 1)))
                width = Inches(float(op.get("width", 2)))
                height = Inches(float(op.get("height", 1)))

                shape = slide.shapes.add_shape(mso_shape, left, top, width, height)

                # Text
                if op.get("text"):
                    shape.text = op["text"]
                    if shape.text_frame.paragraphs:
                        for para in shape.text_frame.paragraphs:
                            para.alignment = align_map.get(op.get("text_alignment", "center"), PP_ALIGN.CENTER)
                            for run in para.runs:
                                _apply_run_format(run, {
                                    "font_size": op.get("font_size"),
                                    "bold": op.get("bold"),
                                    "color": op.get("font_color") or op.get("text_color"),
                                    "font_name": op.get("font_name"),
                                })
                    shape.text_frame.word_wrap = True

                # Fill
                fill_color = op.get("fill_color") or op.get("fill")
                if fill_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = _parse_rgb(fill_color)
                elif op.get("no_fill"):
                    shape.fill.background()

                # Line
                line_color = op.get("line_color")
                if line_color:
                    shape.line.color.rgb = _parse_rgb(line_color)
                if op.get("line_width"):
                    shape.line.width = Pt(float(op["line_width"]))

                # Rotation
                if op.get("rotation"):
                    shape.rotation = float(op["rotation"])

            elif ot == "set_shape_format":
                slide = _get_slide(prs, op)
                shape = _find_shape(slide, op)

                fill_color = op.get("fill_color") or op.get("fill")
                if fill_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = _parse_rgb(fill_color)
                if op.get("no_fill"):
                    shape.fill.background()
                if op.get("line_color"):
                    shape.line.color.rgb = _parse_rgb(op["line_color"])
                if op.get("line_width"):
                    shape.line.width = Pt(float(op["line_width"]))
                if op.get("rotation") is not None:
                    shape.rotation = float(op["rotation"])
                # Resize and reposition
                if op.get("width") is not None:
                    shape.width = Inches(float(op["width"]))
                if op.get("height") is not None:
                    shape.height = Inches(float(op["height"]))
                if op.get("left") is not None:
                    shape.left = Inches(float(op["left"]))
                if op.get("top") is not None:
                    shape.top = Inches(float(op["top"]))

            elif ot == "remove_shape":
                slide = _get_slide(prs, op)
                shape = _find_shape(slide, op)
                sp = shape._element
                sp.getparent().remove(sp)

            elif ot == "add_group_shape":
                slide = _get_slide(prs, op)
                group = slide.shapes.add_group_shape()
                if op.get("name"):
                    group.name = op["name"]

            elif ot == "add_freeform":
                slide = _get_slide(prs, op)
                points = op.get("points", [])
                if len(points) < 2:
                    errors.append(f"Op #{idx} add_freeform: need at least 2 points")
                    continue
                # Accept both [x, y] arrays and {"x": ..., "y": ...} dicts
                def _pt(p):
                    if isinstance(p, (list, tuple)):
                        return float(p[0]), float(p[1])
                    return float(p.get("x", 0)), float(p.get("y", 0))
                sx, sy = _pt(points[0])
                builder = slide.shapes.build_freeform(Inches(sx), Inches(sy))
                for pt in points[1:]:
                    px, py = _pt(pt)
                    builder.add_line_segment(Inches(px), Inches(py))
                if op.get("close", True):
                    builder.add_line_segment(Inches(sx), Inches(sy))
                shape = builder.convert_to_shape()
                fill = op.get("fill_color")
                if fill:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = _parse_rgb(fill)
                if op.get("line_color"):
                    shape.line.color.rgb = _parse_rgb(op["line_color"])

            # =============================================================
            # CONNECTORS
            # =============================================================

            elif ot == "add_connector":
                slide = _get_slide(prs, op)
                from pptx.enum.shapes import MSO_CONNECTOR_TYPE
                start_x = Inches(float(op.get("start_x", 1)))
                start_y = Inches(float(op.get("start_y", 1)))
                end_x = Inches(float(op.get("end_x", 5)))
                end_y = Inches(float(op.get("end_y", 1)))

                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR_TYPE.STRAIGHT, start_x, start_y, end_x, end_y
                )
                if op.get("line_color"):
                    connector.line.color.rgb = _parse_rgb(op["line_color"])
                if op.get("line_width"):
                    connector.line.width = Pt(float(op["line_width"]))

                # Arrow heads via XML
                if op.get("arrow_end") or op.get("arrow_start"):
                    ln = connector.line._ln
                    if op.get("arrow_end"):
                        tail = parse_xml(f'<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="arrow"/>')
                        ln.append(tail)
                    if op.get("arrow_start"):
                        head = parse_xml(f'<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="arrow"/>')
                        ln.append(head)

            # =============================================================
            # IMAGES
            # =============================================================

            elif ot == "add_image":
                slide = _get_slide(prs, op)
                img_path = _resolve_path(
                    op.get("image_path") or op.get("path") or op.get("image", "")
                )
                left = Inches(float(op.get("left", 1)))
                top = Inches(float(op.get("top", 1)))
                width = Inches(float(op["width"])) if op.get("width") else None
                height = Inches(float(op["height"])) if op.get("height") else None
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)

            # =============================================================
            # TABLES
            # =============================================================

            elif ot == "add_table":
                slide = _get_slide(prs, op)
                headers = op.get("headers", [])
                rows = op.get("rows", [])
                cols = len(headers)
                total_rows = 1 + len(rows)
                left = Inches(float(op.get("left", 1)))
                top = Inches(float(op.get("top", 3)))
                width = Inches(float(op.get("width", 8)))
                height = Inches(float(op.get("height", 2)))
                table_shape = slide.shapes.add_table(total_rows, cols, left, top, width, height)
                table = table_shape.table
                for ci, h in enumerate(headers):
                    table.cell(0, ci).text = str(h)
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        if ci < cols:
                            table.cell(ri + 1, ci).text = str(val)

            elif ot == "format_table_cell":
                slide = _get_slide(prs, op)
                ti = int(op.get("table_index", 0))
                tables = [s for s in slide.shapes if s.has_table]
                if ti >= len(tables):
                    errors.append(f"Op #{idx} format_table_cell: table_index {ti} out of range")
                    continue
                table = tables[ti].table
                row = int(op.get("row", 0))
                col = int(op.get("col", 0))

                # Support range
                end_row = int(op.get("end_row", row))
                end_col = int(op.get("end_col", col))

                for r in range(row, end_row + 1):
                    for c in range(col, end_col + 1):
                        if r >= len(table.rows) or c >= len(table.columns):
                            continue
                        cell = table.cell(r, c)
                        # Fill
                        fill_color = op.get("fill_color") or op.get("fill")
                        if fill_color:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _parse_rgb(fill_color)
                        # Text formatting
                        for para in cell.text_frame.paragraphs:
                            if op.get("alignment") and op["alignment"] in align_map:
                                para.alignment = align_map[op["alignment"]]
                            for run in para.runs:
                                _apply_run_format(run, op)
                        # Vertical alignment
                        if op.get("vertical_alignment"):
                            va_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
                            va = va_map.get(op["vertical_alignment"])
                            if va:
                                cell.vertical_anchor = va

                # Merge
                merge_to = op.get("merge_to")
                if merge_to:
                    mr, mc = int(merge_to.get("row", row)), int(merge_to.get("col", col))
                    table.cell(row, col).merge(table.cell(mr, mc))

            # =============================================================
            # CHARTS
            # =============================================================

            elif ot == "add_chart":
                slide = _get_slide(prs, op)
                chart_type_name = op.get("chart_type", "column_clustered")
                ct_map = _get_chart_type_map()
                xl_ct = ct_map.get(chart_type_name)
                if not xl_ct:
                    errors.append(f"Op #{idx} add_chart: unknown chart_type '{chart_type_name}'")
                    continue

                categories = op.get("categories", [])
                series_list = op.get("series", [])

                left = Inches(float(op.get("left", 1)))
                top = Inches(float(op.get("top", 2)))
                width = Inches(float(op.get("width", 8)))
                height = Inches(float(op.get("height", 4.5)))

                # Scatter/XY charts use XyChartData
                if chart_type_name.startswith("scatter"):
                    chart_data = XyChartData()
                    for s in series_list:
                        xy_series = chart_data.add_series(s.get("name", "Series"))
                        x_vals = s.get("x", s.get("values", []))
                        y_vals = s.get("y", [])
                        for xi, yi in zip(x_vals, y_vals):
                            xy_series.add_data_point(xi, yi)
                else:
                    chart_data = CategoryChartData()
                    chart_data.categories = categories
                    for s in series_list:
                        chart_data.add_series(s.get("name", "Series"), s.get("values", []))

                chart_frame = slide.shapes.add_chart(xl_ct, left, top, width, height, chart_data)
                chart = chart_frame.chart

                if op.get("title"):
                    chart.has_title = True
                    chart.chart_title.text_frame.text = op["title"]
                if op.get("has_legend") is not None:
                    chart.has_legend = op["has_legend"]

                # Colors: per-series "color" (hex); pie/doughnut take an
                # op-level "colors" array applied per slice. Without these,
                # decks keep Office default steel-blue amid custom-styled
                # slides.
                try:
                    if chart_type_name in ("pie", "pie_exploded", "doughnut"):
                        slice_colors = op.get("colors") or (
                            series_list[0].get("colors") if series_list else []) or []
                        if slice_colors:
                            for point, col in zip(chart.series[0].points, slice_colors):
                                point.format.fill.solid()
                                point.format.fill.fore_color.rgb = _parse_rgb(col)
                    else:
                        for s_def, s_obj in zip(series_list, chart.series):
                            col = s_def.get("color")
                            if not col:
                                continue
                            rgb = _parse_rgb(col)
                            if chart_type_name.startswith(("line", "scatter", "radar")):
                                s_obj.format.line.color.rgb = rgb
                                try:
                                    s_obj.marker.format.fill.solid()
                                    s_obj.marker.format.fill.fore_color.rgb = rgb
                                except Exception:
                                    pass  # marker styling is best-effort
                            else:
                                s_obj.format.fill.solid()
                                s_obj.format.fill.fore_color.rgb = rgb
                except Exception as exc:
                    errors.append(f"Op #{idx} add_chart: color styling failed: {exc}")

            elif ot == "update_chart_data":
                slide = _get_slide(prs, op)
                # Find chart shape
                ci = int(op.get("chart_index", 0))
                charts = [s for s in slide.shapes if s.has_chart]
                if ci >= len(charts):
                    errors.append(f"Op #{idx} update_chart_data: chart_index {ci} out of range")
                    continue
                chart = charts[ci].chart
                chart_data = CategoryChartData()
                chart_data.categories = op.get("categories", [])
                for s in op.get("series", []):
                    chart_data.add_series(s.get("name", "Series"), s.get("values", []))
                chart.replace_data(chart_data)

            # =============================================================
            # BACKGROUND & TRANSITIONS
            # =============================================================

            elif ot in ("set_background", "set_background_color"):
                slide = _get_slide(prs, op)
                bg = slide.background
                gradient = op.get("gradient")
                if gradient and isinstance(gradient, list) and len(gradient) >= 2:
                    # Build gradient XML directly (python-pptx gradient API is limited)
                    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    stops_xml = ""
                    for i, g in enumerate(gradient):
                        color = _parse_rgb(g.get("color", "FFFFFF"))
                        pos = float(g.get("position", i / (len(gradient) - 1)))
                        pos_pct = int(pos * 100000)
                        stops_xml += (
                            f'<a:gs pos="{pos_pct}" xmlns:a="{ns_a}">'
                            f'<a:srgbClr val="{color}"/>'
                            f'</a:gs>'
                        )
                    angle = int(float(op.get("angle", 270)) * 60000)  # EMU degrees
                    grad_xml = (
                        f'<a:gradFill xmlns:a="{ns_a}">'
                        f'<a:gsLst>{stops_xml}</a:gsLst>'
                        f'<a:lin ang="{angle}" scaled="1"/>'
                        f'</a:gradFill>'
                    )
                    bg_pr = bg._element  # p:bg element
                    # Clear existing fill
                    bgPr = bg_pr.find(qn("p:bgPr"))
                    if bgPr is None:
                        bgPr = parse_xml(
                            f'<p:bgPr xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                            f'<a:noFill xmlns:a="{ns_a}"/>'
                            f'<a:effectLst xmlns:a="{ns_a}"/>'
                            f'</p:bgPr>'
                        )
                        bg_pr.append(bgPr)
                    # Remove old fill children
                    for child in list(bgPr):
                        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        if tag in ("solidFill", "gradFill", "noFill", "pattFill"):
                            bgPr.remove(child)
                    bgPr.insert(0, parse_xml(grad_xml))
                else:
                    fill = bg.fill
                    fill.solid()
                    c = op.get("color", "FFFFFF")
                    fill.fore_color.rgb = _parse_rgb(c)

            elif ot == "set_transition":
                slide = _get_slide(prs, op)
                # Accept multiple key names: transition/effect/name
                transition = op.get("transition") or op.get("effect") or op.get("name") or "fade"
                # If transition is a dict (LLM sends {"type": "fade", "duration": 1000}), extract
                if isinstance(transition, dict):
                    duration = int(transition.get("duration", op.get("duration", 1000)))
                    advance_after = transition.get("advance_after", op.get("advance_after"))
                    transition = transition.get("type") or transition.get("name") or transition.get("effect") or "fade"
                else:
                    duration = int(op.get("duration", 1000))
                    advance_after = op.get("advance_after")

                # XML-based transition
                trans_map = {
                    "fade": "fade", "push": "push", "wipe": "wipe",
                    "split": "split", "reveal": "reveal", "cover": "cover",
                    "dissolve": "dissolve",
                }
                trans_tag = trans_map.get(transition, "fade")
                ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
                ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

                # Remove existing transition
                existing = slide._element.findall(qn("p:transition"))
                for e in existing:
                    slide._element.remove(e)

                trans_el = parse_xml(
                    f'<p:transition xmlns:p="{ns}" spd="med" advTm="{duration}">'
                    f'<p:{trans_tag}/>'
                    f'</p:transition>'
                )
                if advance_after:
                    trans_el.set("advTm", str(advance_after))
                    trans_el.set("advClick", "0")

                slide._element.append(trans_el)

            # =============================================================
            # SPEAKER NOTES
            # =============================================================

            elif ot == "set_speaker_notes":
                slide = _get_slide(prs, op)
                text = op.get("text", "")
                notes_slide = slide.notes_slide
                tf = notes_slide.notes_text_frame
                if op.get("append") and tf.text.strip():
                    tf.text = tf.text + "\n" + text
                else:
                    tf.text = text

            # =============================================================
            # HYPERLINKS
            # =============================================================

            elif ot == "add_hyperlink":
                slide = _get_slide(prs, op)
                url = op.get("url", "")
                shape = _find_shape(slide, op)
                if shape.has_text_frame:
                    # Add hyperlink on a run
                    text = op.get("text")
                    if text:
                        p = shape.text_frame.paragraphs[-1]
                        run = p.add_run()
                        run.text = text
                        run.hyperlink.address = url
                        run.font.color.rgb = _parse_rgb(op.get("color", "0563C1"))
                        run.font.underline = True
                    else:
                        # Hyperlink on whole shape
                        shape.click_action.hyperlink.address = url
                else:
                    shape.click_action.hyperlink.address = url

            # =============================================================
            # SLIDE NUMBERS & PROPERTIES
            # =============================================================

            elif ot == "set_slide_number":
                # Enable slide numbers by adding placeholder to slide master
                enabled = op.get("enabled", True)
                for slide_master in prs.slide_masters:
                    for layout in slide_master.slide_layouts:
                        for ph in layout.placeholders:
                            if ph.placeholder_format.idx == 12:  # Slide number
                                if not enabled:
                                    ph._element.getparent().remove(ph._element)
                                break

            elif ot == "set_core_properties":
                props = prs.core_properties
                if op.get("title"):
                    props.title = op["title"]
                if op.get("author"):
                    props.author = op["author"]
                if op.get("subject"):
                    props.subject = op["subject"]
                if op.get("keywords"):
                    props.keywords = op["keywords"]
                if op.get("comments"):
                    props.comments = op["comments"]

            # =============================================================
            # UNKNOWN
            # =============================================================

            else:
                logger.warning(f"write_pptx: unknown operation '{ot}', skipping")
                errors.append(f"Op #{idx}: unknown operation '{ot}'")

        except Exception as exc:
            errors.append(f"Op #{idx} {ot}: {exc}")
            logger.warning(f"write_pptx op #{idx} '{ot}' failed: {exc}")

    # Save even with partial success
    prs.save(path)
    await _push_preview(path)

    msg = f"Presentation saved: {_to_agents_relative(path)} ({len(ops)} operations applied)"
    if errors:
        msg += f"\n\nWarnings/Errors ({len(errors)}):\n" + "\n".join(f"  - {e}" for e in errors)
    return msg
